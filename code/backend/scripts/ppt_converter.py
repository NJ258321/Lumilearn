#!/usr/bin/env python3
"""
PPT 转图片脚本
纯 Python 实现，无需外部依赖（除了 Python 包）

安装依赖:
pip install python-pptx pymupdf Pillow

如果使用 LibreOffice 会更清晰，但不强制要求
"""

import sys
import os
import json
import tempfile
import shutil
from pathlib import Path

def convert_ppt_to_images(ppt_path, output_dir):
    """
    将 PPT 转换为图片
    策略：
    1. 优先使用 LibreOffice: PPT -> PDF -> 图片
    2. 备用: 使用 python-pptx 直接提取/渲染
    """
    ppt_path = Path(ppt_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    if not ppt_path.exists():
        return {"success": False, "error": "PPT file not found"}

    # 方法1：尝试使用 LibreOffice（效果最好）
    libreoffice_result = try_libreoffice(ppt_path, output_dir)
    if libreoffice_result:
        return libreoffice_result
    
    # 方法2：使用 python-pptx 直接处理
    print("LibreOffice not available, using python-pptx fallback...", file=sys.stderr)
    return convert_using_pptx(ppt_path, output_dir)


def try_libreoffice(ppt_path, output_dir):
    """尝试使用 LibreOffice 转换"""
    try:
        import subprocess
        
        # 查找 LibreOffice
        libreoffice_cmd = None
        for cmd in ['libreoffice', 'soffice', 'soffice.exe']:
            if shutil.which(cmd):
                libreoffice_cmd = cmd
                break
        
        if not libreoffice_cmd:
            return None
        
        print(f"Using {libreoffice_cmd} to convert...", file=sys.stderr)
        
        # PPT -> PDF
        result = subprocess.run([
            libreoffice_cmd, '--headless', '--convert-to', 'pdf',
            '--outdir', str(output_dir), str(ppt_path)
        ], capture_output=True, text=True, timeout=120)
        
        if result.returncode != 0:
            print(f"LibreOffice error: {result.stderr}", file=sys.stderr)
            return None
        
        # 查找生成的 PDF
        pdf_files = list(output_dir.glob("*.pdf"))
        if not pdf_files:
            return None
        
        pdf_path = pdf_files[0]
        print(f"PDF created: {pdf_path}", file=sys.stderr)
        
        # PDF -> 图片 (使用 PyMuPDF，无需 poppler)
        return pdf_to_images(pdf_path, output_dir)
        
    except Exception as e:
        print(f"LibreOffice conversion failed: {e}", file=sys.stderr)
        return None


def pdf_to_images(pdf_path, output_dir):
    """使用 PyMuPDF 将 PDF 转为图片（无需 poppler）"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("PyMuPDF not installed, cannot convert PDF", file=sys.stderr)
        return None
    
    try:
        doc = fitz.open(pdf_path)
        image_paths = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            # 渲染页面 (150 DPI)
            mat = fitz.Matrix(150/72, 150/72)  # 150 DPI
            pix = page.get_pixmap(matrix=mat)
            
            # 保存图片
            img_path = output_dir / f"slide_{page_num + 1}.jpg"
            pix.save(str(img_path))
            image_paths.append(str(img_path))
        
        doc.close()
        
        # 删除临时 PDF
        pdf_path.unlink()
        
        return {
            "success": True,
            "page_count": len(image_paths),
            "images": image_paths
        }
        
    except Exception as e:
        print(f"PDF to images failed: {e}", file=sys.stderr)
        return None


def convert_using_pptx(ppt_path, output_dir):
    """
    使用 python-pptx 直接转换（备用方案）
    提取 PPT 中的图片 + 生成文本预览
    """
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as e:
        return {"success": False, "error": f"Missing package: {e}"}

    try:
        prs = Presentation(ppt_path)
        image_paths = []
        
        # 尝试加载字体
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        except:
            try:
                font = ImageFont.truetype("arial.ttf", 24)
            except:
                font = ImageFont.load_default()
        
        for i, slide in enumerate(prs.slides):
            # 创建白色背景图片 (16:9)
            img = Image.new('RGB', (1280, 720), 'white')
            draw = ImageDraw.Draw(img)
            
            # 提取并绘制文本
            y_offset = 50
            has_content = False
            
            for shape in slide.shapes:
                # 尝试获取文本
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()[:200]  # 限制长度
                    draw.text((50, y_offset), text, fill='black', font=font)
                    y_offset += 35
                    has_content = True
                    
                    if y_offset > 650:  # 防止超出边界
                        break
            
            # 如果没有文本，绘制页码
            if not has_content:
                draw.text((600, 350), f"Slide {i+1}", fill='gray', font=font)
            
            # 保存图片
            img_path = output_dir / f"slide_{i+1}.jpg"
            img.save(img_path, "JPEG", quality=85)
            image_paths.append(str(img_path))
        
        return {
            "success": True,
            "page_count": len(image_paths),
            "images": image_paths,
            "note": "Fallback mode - text only rendering"
        }
        
    except Exception as e:
        return {"success": False, "error": str(e)}


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python ppt_converter.py <ppt_file> [output_dir]", file=sys.stderr)
        sys.exit(1)

    ppt_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else tempfile.mkdtemp()

    result = convert_ppt_to_images(ppt_file, output_dir)
    print(json.dumps(result))
