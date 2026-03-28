#!/usr/bin/env python3
"""
PPT 文本提取脚本
提取每页的标题、正文内容，返回结构化文本

安装依赖:
pip install python-pptx
"""

import sys
import json
from pathlib import Path
from pptx import Presentation


def extract_ppt_content(ppt_path):
    """
    提取 PPT 的文本内容
    返回每页的标题、正文、备注
    """
    ppt_path = Path(ppt_path)
    
    if not ppt_path.exists():
        return {"success": False, "error": "PPT file not found"}
    
    try:
        prs = Presentation(ppt_path)
        slides = []
        
        for i, slide in enumerate(prs.slides):
            slide_data = {
                "page": i + 1,
                "title": "",
                "content": [],
                "notes": ""
            }
            
            # 提取标题（通常是第一个占位符）
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text.strip()
            
            # 提取所有文本内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # 避免重复添加标题
                    if text != slide_data["title"]:
                        slide_data["content"].append(text)
            
            # 提取演讲者备注
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()
            
            # 合并所有内容
            slide_data["full_text"] = f"{slide_data['title']}\n" + "\n".join(slide_data["content"])
            
            slides.append(slide_data)
        
        # 生成 Markdown 格式摘要
        markdown = generate_markdown(slides)
        
        return {
            "success": True,
            "page_count": len(slides),
            "slides": slides,
            "markdown": markdown
        }
        
    except Exception as e:
        return {"success": False, "error": str(e)}


def generate_markdown(slides):
    """生成 Markdown 格式的 PPT 摘要"""
    lines = ["# PPT 内容摘要\n"]
    
    for slide in slides:
        title = slide["title"] or f"第{slide['page']}页"
        lines.append(f"## {title}")
        
        for text in slide["content"]:
            # 如果文本较长，可能是段落，否则作为列表项
            if len(text) > 100:
                lines.append(f"\n{text}")
            else:
                lines.append(f"- {text}")
        
        if slide["notes"]:
            lines.append(f"\n> 备注: {slide['notes']}")
        
        lines.append("")  # 空行
    
    return "\n".join(lines)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python ppt_extractor.py <ppt_file>", file=sys.stderr)
        sys.exit(1)
    
    ppt_file = sys.argv[1]
    result = extract_ppt_content(ppt_file)
    print(json.dumps(result, ensure_ascii=False))
